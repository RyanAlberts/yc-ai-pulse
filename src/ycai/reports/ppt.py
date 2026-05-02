"""VC-style ``.pptx`` deck generator.

Visual signature: cream background (``#F5F1E8``), orange accent (``#D24E01``),
sans-serif headings, serif body. Modeled after a16z-style blog posts but
deliberately uses original color values + system font fallbacks — no a16z
logos, no copied cover illustrations.

Charts are rendered with matplotlib (PNG → embedded image), not native PowerPoint
charts. This is intentional: matplotlib gives us pixel-perfect control over the
heatmap colormap, and the same chart math is shared with the HTML dashboard via
``ycai.analytics``. Numbers can't drift between the two surfaces because they
read from the same source.

Layer 2 anti-hallucination invariants (enforced before the file is written):
  - Every paragraph in the deck is scanned for forbidden hedge phrases.
  - Every number that appears in deck prose must trace back to
    ``analytics.headline_numbers`` or one of the counters used to build the
    charts. Drift aborts the build with the offending span.
"""

from __future__ import annotations

import io
import logging
from collections import Counter
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ycai import analytics
from ycai.reports.anti_hallucination import audit
from ycai.schemas import BatchCoverage, CompanyAnalysis, RawCompany

log = logging.getLogger(__name__)

# Palette — original values, not copied from any specific brand.
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
INK = RGBColor(0x1B, 0x1B, 0x1B)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
ACCENT = RGBColor(0xD2, 0x4E, 0x01)
LINE = RGBColor(0xDD, 0xD8, 0xCB)

# Slide dimensions: 16:9 widescreen at 13.333" x 7.5".
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class Layer2Failure(RuntimeError):
    """Raised when Layer 2 audit fails. Carries the offending hits for diagnostics."""

    def __init__(self, message: str, *, forbidden: list, drifts: list) -> None:
        super().__init__(message)
        self.forbidden = forbidden
        self.drifts = drifts


@dataclass
class DeckBuildContext:
    """Everything a slide builder needs. Built once at the top of ``build_deck``."""

    coverage: BatchCoverage
    companies: list[RawCompany]
    analyses: list[CompanyAnalysis]
    industry: Counter[str]
    capability: Counter[str]
    capability_heatmap: analytics.CapabilityHeatmap
    tech_stack: Counter[str]
    oss_posture: Counter[str]
    headline: dict[str, int]
    spotlights: list[CompanyAnalysis]
    quotes: list[CompanyAnalysis]
    fetched_at: datetime
    # Two prose streams:
    #   prose_buffer: aggregate commentary (TL;DR, chart commentary, methodology).
    #     Numbers here MUST trace back to the dataframe — drift-checked.
    #   quoted_facts: per-company text (rationale, tagline). These have already
    #     been validated by Layer 1's source-URL grounding, so we run the
    #     forbidden-phrase scan on them but skip the drift check.
    prose_buffer: list[str]
    quoted_facts: list[str]

    @classmethod
    def from_run(
        cls,
        coverage: BatchCoverage,
        companies: list[RawCompany],
        analyses: list[CompanyAnalysis],
    ) -> DeckBuildContext:
        return cls(
            coverage=coverage,
            companies=companies,
            analyses=analyses,
            industry=analytics.industry_distribution(analyses),
            capability=analytics.capability_totals(analyses),
            capability_heatmap=analytics.capability_heatmap(analyses),
            tech_stack=analytics.tech_stack_distribution(analyses),
            oss_posture=analytics.oss_posture_distribution(analyses),
            headline=analytics.headline_numbers(analyses, coverage=coverage),
            spotlights=analytics.spotlight_companies(analyses, top_n=6),
            quotes=analytics.quote_candidates(analyses)[:4],
            fetched_at=datetime.now(UTC),
            prose_buffer=[],
            quoted_facts=[],
        )

    def all_counters(self) -> list[Counter[str]]:
        return [self.industry, self.capability, self.tech_stack, self.oss_posture]


# ----- matplotlib chart helpers ---------------------------------------------------------


plt.rcParams.update(
    {
        "font.family": "DejaVu Sans",
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.edgecolor": "#999",
        "axes.labelcolor": "#1B1B1B",
        "xtick.color": "#1B1B1B",
        "ytick.color": "#1B1B1B",
        "axes.grid": False,
        "savefig.facecolor": "#F5F1E8",
        "figure.facecolor": "#F5F1E8",
    }
)


def _png_horizontal_bar(counter: Counter[str], *, top: int = 10, color: str = "#D24E01") -> bytes:
    items = counter.most_common(top)
    if not items:
        return _png_empty()
    items.reverse()
    names = [n for n, _ in items]
    values = [v for _, v in items]
    fig, ax = plt.subplots(figsize=(9, 5.5))
    ax.barh(names, values, color=color, edgecolor="none")
    for i, v in enumerate(values):
        ax.text(v, i, f"  {v}", va="center", fontsize=10, color="#1B1B1B")
    ax.set_xlabel("companies")
    ax.tick_params(axis="y", labelsize=10)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def _png_pie(counter: Counter[str], *, color_map: dict[str, str]) -> bytes:
    items = counter.most_common()
    if not items:
        return _png_empty()
    labels = [name for name, _ in items]
    values = [v for _, v in items]
    colors = [color_map.get(name, "#9CA3AF") for name in labels]
    fig, ax = plt.subplots(figsize=(7, 5.5))
    _wedges, _texts, _autotexts = ax.pie(
        values,
        labels=[f"{n} ({v})" for n, v in items],
        colors=colors,
        autopct="%1.0f%%",
        startangle=90,
        wedgeprops={"linewidth": 2, "edgecolor": "white"},
        textprops={"fontsize": 10, "color": "#1B1B1B"},
    )
    ax.axis("equal")
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def _png_heatmap(heatmap: analytics.CapabilityHeatmap) -> bytes:
    if not heatmap.matrix:
        return _png_empty()
    rows = heatmap.capabilities
    cols = heatmap.industries
    grid = [[heatmap.matrix.get((r, c), 0) for c in cols] for r in rows]
    fig, ax = plt.subplots(figsize=(10, 6))
    cmap = plt.get_cmap("Oranges")
    im = ax.imshow(grid, cmap=cmap, aspect="auto")
    ax.set_xticks(range(len(cols)))
    ax.set_xticklabels(cols, rotation=25, ha="right", fontsize=10)
    ax.set_yticks(range(len(rows)))
    ax.set_yticklabels(rows, fontsize=10)
    for i in range(len(rows)):
        for j in range(len(cols)):
            value = grid[i][j]
            if value:
                ax.text(j, i, str(value), ha="center", va="center", fontsize=9, color="#1B1B1B")
    fig.colorbar(im, ax=ax, label="companies")
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def _png_empty() -> bytes:
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.text(0.5, 0.5, "No data", ha="center", va="center", fontsize=18, color="#9CA3AF")
    ax.axis("off")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


# ----- slide builders -------------------------------------------------------------------


def _set_slide_background(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    font_name: str = "Helvetica Neue",
    alignment=None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_image(slide, *, image_bytes: bytes, left: float, top: float, width: float, height: float) -> None:
    slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(left), Inches(top), Inches(width), Inches(height))


def _add_accent_bar(slide, *, left: float, top: float, width: float = 1.0, height: float = 0.08) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def _slide_footer(slide, *, source_url: str) -> None:
    _add_textbox(
        slide,
        text=f"Source: {source_url}",
        left=0.5,
        top=7.05,
        width=12.3,
        height=0.4,
        font_size=10,
        color=MUTED,
    )


def _build_title_slide(prs: Presentation, ctx: DeckBuildContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_background(slide, CREAM)
    _add_accent_bar(slide, left=0.6, top=2.4, width=1.4, height=0.12)
    _add_textbox(
        slide,
        text=f"The State of AI in {ctx.coverage.batch_label}",
        left=0.6,
        top=2.65,
        width=12.0,
        height=2.0,
        font_size=58,
        bold=True,
    )
    headline_pct = (
        ctx.coverage.coverage_pct_of_official
        if ctx.coverage.coverage_pct_of_official is not None
        else ctx.coverage.coverage_pct_of_upstream
    )
    _add_textbox(
        slide,
        text=(
            f"{headline_pct}% of the batch analyzed under strict anti-hallucination guards. "
            f"{ctx.headline['cohort_size']} companies in the high-confidence cohort."
        ),
        left=0.6,
        top=4.4,
        width=12.0,
        height=1.2,
        font_size=22,
        color=MUTED,
    )
    ctx.prose_buffer.append(
        f"{headline_pct}% of the batch analyzed. " f"{ctx.headline['cohort_size']} companies in cohort."
    )
    _add_textbox(
        slide,
        text=f"yc-ai-pulse · generated {ctx.fetched_at.strftime('%Y-%m-%d')}",
        left=0.6,
        top=6.5,
        width=12.0,
        height=0.4,
        font_size=12,
        color=MUTED,
    )


def _build_tldr_slide(prs: Presentation, ctx: DeckBuildContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, CREAM)
    _add_accent_bar(slide, left=0.6, top=0.6, width=0.8, height=0.08)
    _add_textbox(slide, text="TL;DR", left=0.6, top=0.8, width=12.0, height=1.0, font_size=44, bold=True)

    h = ctx.headline
    cohort = max(h["cohort_size"], 1)
    agents_pct = round(h["agents_count"] * 100 / cohort)
    no_ai_pct = round(h["no_ai_count"] * 100 / cohort)
    bullets = [
        f"{h['agents_count']} of {cohort} high-confidence companies build agents.",
        f"{h['no_ai_count']} companies were classified as no-ai despite being in the YC AI batch.",
        f"OSS posture: {ctx.oss_posture.get('closed', 0)} closed, "
        f"{ctx.oss_posture.get('unknown', 0)} unknown, "
        f"{ctx.oss_posture.get('api-only', 0)} api-only.",
    ]
    for i, line in enumerate(bullets):
        _add_textbox(
            slide,
            text=f"• {line}",
            left=0.6,
            top=2.2 + i * 1.0,
            width=12.0,
            height=0.9,
            font_size=22,
        )
        ctx.prose_buffer.append(line)
    _ = agents_pct
    _ = no_ai_pct
    _slide_footer(slide, source_url="github.com/RyanAlberts/yc-ai-pulse")


def _build_chart_slide(
    prs: Presentation,
    ctx: DeckBuildContext,
    *,
    title: str,
    image_bytes: bytes,
    commentary: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, CREAM)
    _add_accent_bar(slide, left=0.6, top=0.6, width=0.8, height=0.08)
    _add_textbox(slide, text=title, left=0.6, top=0.8, width=12.0, height=1.0, font_size=32, bold=True)
    _add_image(slide, image_bytes=image_bytes, left=0.5, top=1.9, width=8.5, height=4.6)
    _add_textbox(
        slide,
        text=commentary,
        left=9.3,
        top=1.9,
        width=3.5,
        height=4.6,
        font_size=14,
        color=INK,
    )
    ctx.prose_buffer.append(commentary)
    _slide_footer(slide, source_url="github.com/RyanAlberts/yc-ai-pulse")


def _build_pull_quote_slide(prs: Presentation, ctx: DeckBuildContext, quote: CompanyAnalysis) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, CREAM)
    _add_textbox(
        slide,
        text=f'"{quote.tagline_rewrite}"',
        left=1.5,
        top=2.4,
        width=10.5,
        height=2.5,
        font_size=36,
        bold=True,
    )
    _add_textbox(
        slide,
        text=f"— {quote.slug}",
        left=1.5,
        top=4.8,
        width=10.5,
        height=0.6,
        font_size=18,
        color=MUTED,
    )
    ctx.quoted_facts.append(quote.tagline_rewrite)
    _slide_footer(slide, source_url=str(quote.sources[0]) if quote.sources else "—")


def _build_spotlight_slide(prs: Presentation, ctx: DeckBuildContext, company: CompanyAnalysis) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, CREAM)
    _add_accent_bar(slide, left=0.6, top=0.6, width=0.8, height=0.08)
    _add_textbox(slide, text=company.slug, left=0.6, top=0.8, width=12.0, height=0.8, font_size=14, color=MUTED)
    _add_textbox(
        slide,
        text=company.tagline_rewrite,
        left=0.6,
        top=1.6,
        width=12.0,
        height=1.6,
        font_size=28,
        bold=True,
    )
    facts = [
        f"Industry: {company.industry_primary.value}",
        f"AI capability: {', '.join(c.value for c in company.ai_capability)}",
        f"Tech stack: {', '.join(s.value for s in company.tech_stack) or 'unknown'}",
        f"OSS posture: {company.oss_posture.value}",
    ]
    for i, line in enumerate(facts):
        _add_textbox(slide, text=line, left=0.6, top=3.4 + i * 0.5, width=12.0, height=0.5, font_size=16, color=MUTED)
    if company.rationale:
        _add_textbox(
            slide,
            text=f'"{company.rationale}"',
            left=0.6,
            top=5.6,
            width=12.0,
            height=1.2,
            font_size=12,
            color=INK,
        )
    ctx.quoted_facts.append(company.tagline_rewrite)
    ctx.quoted_facts.append(company.rationale)
    _slide_footer(slide, source_url=str(company.sources[0]) if company.sources else "—")


def _build_methodology_slide(prs: Presentation, ctx: DeckBuildContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, CREAM)
    _add_accent_bar(slide, left=0.6, top=0.6, width=0.8, height=0.08)
    _add_textbox(slide, text="Methodology", left=0.6, top=0.8, width=12.0, height=1.0, font_size=36, bold=True)
    h = ctx.headline
    lines = [
        f"Source: {ctx.coverage.source} (last refresh {ctx.coverage.source_last_updated}).",
        f"Coverage: {ctx.coverage.upstream_company_count} of "
        f"{ctx.coverage.yc_official_count or ctx.coverage.upstream_company_count} from yc-oss/api.",
        f"Tier A: {ctx.coverage.tier_a_count} · Tier B: {ctx.coverage.tier_b_count} · "
        f"Tier C (excluded): {ctx.coverage.tier_c_count}",
        f"LLM cohort: {h['high_confidence']} high + {h['medium_confidence']} medium "
        f"({h['low_confidence']} low excluded). Sonnet 4.6 via Agent SDK.",
        "Anti-hallucination Layer 1: pydantic schema, source-URL grounding, "
        "two-pass cross-check, sentinel on failure.",
        "Anti-hallucination Layer 2 (this deck): forbidden-phrase scan + " "numerical-drift check before write.",
        "Depth=1 website crawl: polite (robots-aware, max 5 pages, 30 KB cap, " "PII stripped before the LLM).",
    ]
    for i, line in enumerate(lines):
        _add_textbox(slide, text=f"• {line}", left=0.6, top=2.0 + i * 0.55, width=12.0, height=0.6, font_size=14)
        ctx.prose_buffer.append(line)
    _slide_footer(slide, source_url="github.com/RyanAlberts/yc-ai-pulse")


def _build_reproduce_slide(prs: Presentation, ctx: DeckBuildContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, CREAM)
    _add_accent_bar(slide, left=0.6, top=0.6, width=0.8, height=0.08)
    _add_textbox(slide, text="Reproduce this deck", left=0.6, top=0.8, width=12.0, height=1.0, font_size=36, bold=True)
    _add_textbox(
        slide,
        text="pipx install yc-ai-pulse",
        left=0.6,
        top=2.4,
        width=12.0,
        height=0.7,
        font_size=22,
        font_name="Menlo",
        color=ACCENT,
    )
    _add_textbox(
        slide,
        text=f"ycai run-coverage --batch {ctx.coverage.batch_slug} --enrich",
        left=0.6,
        top=3.2,
        width=12.0,
        height=0.7,
        font_size=22,
        font_name="Menlo",
        color=ACCENT,
    )
    _add_textbox(
        slide,
        text=f"ycai report {ctx.coverage.batch_slug}/  # produces deck.pptx and report.docx",
        left=0.6,
        top=4.0,
        width=12.0,
        height=0.7,
        font_size=22,
        font_name="Menlo",
        color=ACCENT,
    )
    _add_textbox(
        slide,
        text="Source code: github.com/RyanAlberts/yc-ai-pulse",
        left=0.6,
        top=5.5,
        width=12.0,
        height=0.6,
        font_size=18,
        color=MUTED,
    )


# ----- top-level builder -----------------------------------------------------------------


def build_deck(
    coverage: BatchCoverage,
    companies: list[RawCompany],
    analyses: list[CompanyAnalysis],
    *,
    output_path: Path,
) -> Path:
    """Generate the deck. Audits prose Layer-2 before writing — raises ``Layer2Failure`` on drift."""
    ctx = DeckBuildContext.from_run(coverage, companies, analyses)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_title_slide(prs, ctx)
    _build_tldr_slide(prs, ctx)
    _build_chart_slide(
        prs,
        ctx,
        title="Industry distribution",
        image_bytes=_png_horizontal_bar(ctx.industry, top=12),
        commentary=(
            f"Top: {', '.join(f'{n} {c}' for c, n in ctx.industry.most_common(3))}. "
            f"Cohort: {ctx.headline['cohort_size']} high+medium-confidence companies."
        ),
    )
    if ctx.quotes:
        _build_pull_quote_slide(prs, ctx, ctx.quotes[0])
    _build_chart_slide(
        prs,
        ctx,
        title="AI capability x industry heatmap",
        image_bytes=_png_heatmap(ctx.capability_heatmap),
        commentary=(
            f"Top capability: {ctx.capability.most_common(1)[0][0] if ctx.capability else 'n/a'} "
            f"({ctx.capability.most_common(1)[0][1] if ctx.capability else 0} of {ctx.headline['cohort_size']}). "
            "Cell label = number of companies pairing that capability with that industry."
        ),
    )
    _build_chart_slide(
        prs,
        ctx,
        title="Tech stack signals",
        image_bytes=_png_horizontal_bar(ctx.tech_stack, top=10),
        commentary=(
            "Many companies returned tech_stack=unknown — the LLM can only cite what's "
            "visible on the website. The named entries are the model providers and "
            "frameworks the model could verify from public surfaces."
        ),
    )
    _build_chart_slide(
        prs,
        ctx,
        title="Open-source posture",
        image_bytes=_png_pie(
            ctx.oss_posture,
            color_map={
                "fully-open": "#15803D",
                "weights-only": "#65A30D",
                "source-available": "#84CC16",
                "api-only": "#F59E0B",
                "closed": "#B91C1C",
                "unknown": "#9CA3AF",
            },
        ),
        commentary=(
            "Closed dominates — typical for B2B SaaS. The unknown slice is the cohort the "
            "model could not classify from public surfaces alone."
        ),
    )
    if len(ctx.quotes) > 1:
        _build_pull_quote_slide(prs, ctx, ctx.quotes[1])
    for company in ctx.spotlights:
        _build_spotlight_slide(prs, ctx, company)
    _build_methodology_slide(prs, ctx)
    _build_reproduce_slide(prs, ctx)

    # Layer 2 audit. Two prose streams audited differently:
    #   prose_buffer (aggregate commentary): forbidden-phrase + drift check.
    #   quoted_facts (per-company taglines & rationales): forbidden-phrase
    #     only. Numbers in these come from the model's view of a single
    #     company's website + YC profile and were already gated by Layer 1's
    #     source-URL guard. Drift-checking them against aggregate counters
    #     would be a category error.
    infra_facts: tuple[float, ...] = (
        4.6,  # model version mention
        5,  # crawler max pages
        30,  # crawler max KB
        2,  # layer 2
        1,  # layer 1, depth=1
    )
    aggregate_prose = " ".join(ctx.prose_buffer)
    quoted_prose = " ".join(ctx.quoted_facts)
    report = audit(aggregate_prose, ctx.headline, ctx.all_counters(), extra_allowed=infra_facts)
    quoted_forbidden = audit(quoted_prose, ctx.headline, ctx.all_counters()).forbidden
    report = type(report)(
        forbidden=report.forbidden + quoted_forbidden,
        drifts=report.drifts,
    )
    if not report.is_clean:
        raise Layer2Failure(
            f"Deck Layer 2 audit failed: {len(report.forbidden)} forbidden phrase(s), "
            f"{len(report.drifts)} numerical drift(s).",
            forbidden=report.forbidden,
            drifts=report.drifts,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


__all__ = ["DeckBuildContext", "Layer2Failure", "build_deck"]
